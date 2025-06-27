import os
import pandas as pd
import docx
from openpyxl import load_workbook
from pptx import Presentation
import pinecone
from langchain_community.document_loaders import PyPDFLoader
from langchain_huggingface import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from pinecone import Pinecone, ServerlessSpec

# Folder where documents are stored
DOCUMENTS_FOLDER = 'documents'  # Ensure you have your documents in this folder

# Pinecone Configuration
PINECONE_API_KEY = "pcsk_RtTvf_FprhEhvmE6aJGckZg9P14Rny1q19p1QnGuRt67eCDWLMzJe3yW1LqaJ8rs1RcyE"  # Replace with your Pinecone API Key
PINECONE_ENV = "us-east-1"  # Replace with your Pinecone environment (you can find it on the Pinecone dashboard)
PINECONE_HOST = "https://pdfrag-bxptziy.svc.aped-4627-b74a.pinecone.io"  # The host URL provided by Pinecone
PINECONE_INDEX_NAME = "pdfrag"  # Name of the Pinecone index

# Initialize Pinecone using the new method
# Create a Pinecone client instance
pc = Pinecone(api_key=PINECONE_API_KEY)

# Check and create the Pinecone index with dimension 384 if it doesn't exist
def create_pinecone_index():
    """Create Pinecone index if not exists."""
    if PINECONE_INDEX_NAME not in pc.list_indexes().names():
        # Create a new index with dimension 384 (matching HuggingFace embeddings)
        pc.create_index(
            name=PINECONE_INDEX_NAME,
            dimension=384,  # Set dimension to 384 (as your embeddings are of size 384)
            metric="cosine",
            spec=ServerlessSpec(
                cloud="aws",
                region="us-east-1"
            )
        )
        print(f"Pinecone index '{PINECONE_INDEX_NAME}' created with dimension 384.")
    else:
        print(f"Pinecone index '{PINECONE_INDEX_NAME}' already exists.")

# Initialize HuggingFace Embeddings model
embeddings = HuggingFaceEmbeddings(model_name='sentence-transformers/all-MiniLM-L6-v2')

def load_pdf(file_path: str) -> list[Document]:
    """Load a PDF file and return its documents."""
    try:
        loader = PyPDFLoader(file_path)
        documents = loader.load()
        print(f"Loaded {len(documents)} documents from {file_path}")
        return documents
    except FileNotFoundError:
        print(f"Error: File {file_path} not found.")
        return []
    except Exception as e:
        print(f"Error loading PDF: {e}")
        return []

def load_csv(file_path: str) -> list[Document]:
    """Load a CSV file and return its documents."""
    try:
        df = pd.read_csv(file_path)
        documents = []
        
        # Convert each row to a document
        for index, row in df.iterrows():
            # Create a text representation of the row
            row_text = ", ".join([f"{col}: {str(val)}" for col, val in row.items() if pd.notna(val)])
            doc = Document(
                page_content=row_text,
                metadata={"source": file_path, "row": index, "type": "csv"}
            )
            documents.append(doc)
        
        print(f"Loaded {len(documents)} rows from CSV {file_path}")
        return documents
    except Exception as e:
        print(f"Error loading CSV: {e}")
        return []

def load_docx(file_path: str) -> list[Document]:
    """Load a DOCX file and return its documents."""
    try:
        doc = docx.Document(file_path)
        documents = []
        
        # Extract text from paragraphs
        full_text = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                full_text.append(paragraph.text.strip())
        
        # Create a single document with all text
        if full_text:
            content = "\n".join(full_text)
            document = Document(
                page_content=content,
                metadata={"source": file_path, "type": "docx"}
            )
            documents.append(document)
        
        print(f"Loaded {len(documents)} documents from DOCX {file_path}")
        return documents
    except Exception as e:
        print(f"Error loading DOCX: {e}")
        return []

def load_excel(file_path: str) -> list[Document]:
    """Load an Excel file (XLSX/XLS) and return its documents."""
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        documents = []
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # Convert each row to a document
            for index, row in df.iterrows():
                row_text = ", ".join([f"{col}: {str(val)}" for col, val in row.items() if pd.notna(val)])
                doc = Document(
                    page_content=row_text,
                    metadata={"source": file_path, "sheet": sheet_name, "row": index, "type": "excel"}
                )
                documents.append(doc)
        
        print(f"Loaded {len(documents)} rows from Excel {file_path}")
        return documents
    except Exception as e:
        print(f"Error loading Excel: {e}")
        return []

def load_pptx(file_path: str) -> list[Document]:
    """Load a PPTX file and return its documents."""
    try:
        presentation = Presentation(file_path)
        documents = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # Create a document for each slide if it has content
            if slide_text:
                content = "\n".join(slide_text)
                document = Document(
                    page_content=content,
                    metadata={"source": file_path, "slide": slide_num, "type": "pptx"}
                )
                documents.append(document)
        
        print(f"Loaded {len(documents)} slides from PPTX {file_path}")
        return documents
    except Exception as e:
        print(f"Error loading PPTX: {e}")
        return []

def load_txt(file_path: str) -> list[Document]:
    """Load a TXT file and return its documents."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        if content.strip():
            document = Document(
                page_content=content,
                metadata={"source": file_path, "type": "txt"}
            )
            print(f"Loaded 1 document from TXT {file_path}")
            return [document]
        return []
    except Exception as e:
        print(f"Error loading TXT: {e}")
        return []

def load_document(file_path: str) -> list[Document]:
    """Load a document based on its file extension."""
    file_extension = os.path.splitext(file_path)[1].lower()
    
    if file_extension == '.pdf':
        return load_pdf(file_path)
    elif file_extension == '.csv':
        return load_csv(file_path)
    elif file_extension == '.docx':
        return load_docx(file_path)
    elif file_extension in ['.xlsx', '.xls']:
        return load_excel(file_path)
    elif file_extension == '.txt':
        return load_txt(file_path)
    elif file_extension == '.pptx':
        return load_pptx(file_path)
    else:
        print(f"Unsupported file format: {file_extension} for file {file_path}")
        return []

def split_documents(documents: list[Document], chunk_size: int = 500, chunk_overlap: int = 100) -> list[Document]:
    """Split documents into smaller chunks."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", " ", ""]
    )
    return text_splitter.split_documents(documents)

def insert_documents_into_pinecone(documents: list[Document]):
    """Insert document embeddings into Pinecone."""
    try:
        # Connect to the Pinecone index
        index = pc.Index(PINECONE_INDEX_NAME)

        # Prepare the embeddings and store them in Pinecone
        batch_size = 100  # Define the batch size
        batch = []

        # Get the embeddings for all documents
        document_texts = [doc.page_content for doc in documents]  # Extract text from documents
        embeddings_list = embeddings.embed_documents(document_texts)  # Use embed_documents method

        for i, doc in enumerate(documents):
            embedding = embeddings_list[i]  # Get the embedding for the current document
            metadata = {
                "text": doc.page_content,
                **doc.metadata  # Include original metadata from the document
            }
            # Use document ID (e.g., a simple ID based on the index) to store the embeddings
            batch.append((f"doc_{i}", embedding, metadata))
            
            # Insert embeddings into Pinecone in batches
            if len(batch) >= batch_size:
                index.upsert(vectors=batch)
                print(f"Upserted {len(batch)} vectors into Pinecone.")
                batch = []  # Clear the batch

        # Upsert any remaining documents
        if batch:
            index.upsert(vectors=batch)
            print(f"Upserted {len(batch)} vectors into Pinecone.")

    except Exception as e:
        print(f"Error inserting documents into Pinecone: {e}")

def get_supported_files(folder_path: str) -> list[str]:
    """Get all supported files from the folder."""
    supported_extensions = ['.pdf', '.csv', '.docx', '.xlsx', '.xls', '.txt', '.pptx']
    files = []
    
    for file in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file)
        if os.path.isfile(file_path):
            file_extension = os.path.splitext(file)[1].lower()
            if file_extension in supported_extensions:
                files.append(file_path)
    
    return files

def process_documents_and_create_pinecone():
    """Process all supported documents from the 'documents' folder and create a vector store in Pinecone."""
    all_documents = []

    # Get all supported files from the documents folder
    document_files = get_supported_files(DOCUMENTS_FOLDER)

    if not document_files:
        print(f"No supported files found in the folder '{DOCUMENTS_FOLDER}'.")
        print("Supported formats: PDF, CSV, DOCX, XLSX, XLS, TXT, PPTX")
        return

    print(f"Found {len(document_files)} supported files:")
    for file in document_files:
        print(f"  - {os.path.basename(file)}")

    # Load and process each document
    for document_file in document_files:
        print(f"\nProcessing: {os.path.basename(document_file)}")
        documents = load_document(document_file)
        if not documents:
            print(f"Failed to load document: {document_file}")
            continue
        all_documents.extend(documents)

    if not all_documents:
        print("No documents were successfully loaded.")
        return

    print(f"\nTotal documents loaded: {len(all_documents)}")

    # Split documents into smaller chunks
    split_docs = split_documents(all_documents)
    print(f"Total chunks after splitting: {len(split_docs)}")

    # Create Pinecone index if it doesn't already exist
    create_pinecone_index()

    # Insert documents into Pinecone
    insert_documents_into_pinecone(split_docs)
    
    print(f"\nSuccessfully processed and stored {len(split_docs)} document chunks in Pinecone!")

if __name__ == "__main__":
    process_documents_and_create_pinecone()